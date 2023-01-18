from pptx import Presentation
import openpyxl
wb = openpyxl.load_workbook("Spring23 Roster.xlsx")
ws = wb['Sheet1']
ppt = Presentation('Spring23 Door Decs.pptx')
layout_cah = ppt.slide_layouts[0]

for x in range(1, 63):
    first_name = 'B' + str(x)
    last_name = 'A' + str(x)
    new_slide = ppt.slides.add_slide(layout_cah)
    new_slide.placeholders[13].text  = str(ws[first_name].value + ' ' + ws[last_name].value)

ppt.save("Spring23 Door Decs.pptx")